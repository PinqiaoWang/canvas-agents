"""
Embedding Pipeline for Canvas Course Documents

A class-based system for processing course documents, extracting text,
creating embeddings, and storing them in Pinecone vector database.

Provides document processing, text chunking, embedding generation,
and vector storage with comprehensive logging and error handling.

Author: Vishwanath Guruvayur
"""

import os
import json
import re
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Any
from datetime import datetime
from dotenv import load_dotenv
import torch
from io import BytesIO

from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document
from sentence_transformers import SentenceTransformer
from pinecone import Pinecone, ServerlessSpec

# Import PyMuPDF for fast PDF processing
try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False
    print("Warning: PyMuPDF not available. Install with: pip install PyMuPDF")

# Import python-pptx for PowerPoint processing
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx not available. Install with: pip install python-pptx")

# Import BeautifulSoup for HTML parsing
try:
    from bs4 import BeautifulSoup
    BS4_AVAILABLE = True
except ImportError:
    BS4_AVAILABLE = False
    print("Warning: BeautifulSoup4 not available. Install with: pip install beautifulsoup4 lxml")

# Import PIL for image handling
try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    print("Warning: Pillow not available. Install with: pip install pillow")


class DocumentProcessor:
    """
    Handles document text extraction and processing.
    
    Supports PDF files using Docling and plain text files,
    with robust error handling and text chunking capabilities.
    """
    
    def __init__(self, chunk_size: int = 500, chunk_overlap: int = 100):
        """
        Initialize the document processor.
        
        Args:
            chunk_size: Maximum size of text chunks
            chunk_overlap: Overlap between consecutive chunks
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separators=["\n\n", "\n", ".", " ", ""]
        )
    
    def extract_pdf_content(self, pdf_path: Path) -> Dict[str, Any]:
        """
        Extract text and images from PDF using PyMuPDF (fast and reliable).
        
        Args:
            pdf_path: Path to the PDF file
            
        Returns:
            Dictionary with 'text' (str) and 'images' (List[Tuple[int, Image.Image]])
            where images are (page_number, PIL.Image) tuples
            
        Raises:
            Exception: If PDF processing fails
        """
        if not PYMUPDF_AVAILABLE:
            raise Exception("PyMuPDF not available. Install with: pip install PyMuPDF")
        
        try:
            doc = fitz.open(pdf_path)
            pages = []
            images = []
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                # Use get_text() for fast extraction with basic formatting
                text = page.get_text().strip()
                
                if text and self._is_valid_text(text):
                    pages.append(f"[Page {page_num + 1}]\n{text}")
                else:
                    print(f"Warning: Page {page_num + 1} contains invalid/empty text, skipping")
                
                # Extract images from this page if PIL is available
                if PIL_AVAILABLE:
                    try:
                        for img_index, img in enumerate(page.get_images(full=True)):
                            xref = img[0]
                            try:
                                base_image = doc.extract_image(xref)
                                image_bytes = base_image.get("image")
                                if image_bytes:
                                    pil_image = Image.open(BytesIO(image_bytes)).convert("RGB")
                                    images.append((page_num + 1, pil_image))
                                    print(f"Extracted image {img_index + 1} from page {page_num + 1}")
                            except Exception as e:
                                print(f"Warning: Could not extract image {img_index + 1} from page {page_num + 1}: {e}")
                                continue
                    except Exception as e:
                        print(f"Warning: Error extracting images from page {page_num + 1}: {e}")
            
            doc.close()
            
            if not pages:
                raise ValueError(f"No valid text could be extracted from PDF {pdf_path.name}")
            
            text_content = "\n".join(pages).strip()
            print(f"Successfully extracted text from {len(pages)} pages and {len(images)} images from {pdf_path.name}")
            
            return {
                "text": text_content,
                "images": images
            }

        except Exception as e:
            raise Exception(f"Failed to process PDF {pdf_path.name}: {e}")
    
    def _is_valid_text(self, text: str) -> bool:
        """
        Check if extracted text is valid (not PDF binary artifacts).
        
        Args:
            text: Text to validate
            
        Returns:
            True if text appears to be valid content
        """
        # Check for common PDF binary artifacts
        invalid_patterns = [
            '%PDF-',
            'obj',
            'endobj',
            'stream',
            'endstream',
            'xref',
            'trailer',
            '/Type/',
            '/Filter/',
            '/Length',
            '<<',
            '>>',
            'R ',
        ]
        
        # If text is mostly binary artifacts, it's invalid
        invalid_count = sum(1 for pattern in invalid_patterns if pattern in text)
        if invalid_count > 3:
            return False
        
        # Check if text has reasonable character distribution
        if len(text) < 10:
            return False
            
        # Check for reasonable ratio of printable characters
        printable_chars = sum(1 for c in text if c.isprintable() or c.isspace())
        if printable_chars / len(text) < 0.7:
            return False
            
        return True
    
    def extract_text_file(self, text_path: Path) -> str:
        """
        Extract text from plain text file.
        
        Args:
            text_path: Path to the text file
            
        Returns:
            File contents as string
            
        Raises:
            Exception: If file reading fails
        """
        try:
            with open(text_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read().strip()
                if not content:
                    raise ValueError(f"File {text_path.name} is empty")
                return content
                
        except Exception as e:
            raise Exception(f"Failed to read text file {text_path.name}: {e}")
    
    def extract_ipynb_text(self, ipynb_path: Path) -> str:
        """
        Extract text from Jupyter notebook (.ipynb) file.
        
        Args:
            ipynb_path: Path to the .ipynb file
            
        Returns:
            Extracted text content from code cells and markdown cells
            
        Raises:
            Exception: If notebook processing fails
        """
        print("IPYDF")
        try:
            with open(ipynb_path, "r", encoding="utf-8") as f:
                notebook_data = json.load(f)
            
            if "cells" not in notebook_data:
                raise ValueError(f"Invalid notebook format in {ipynb_path.name}")
            
            extracted_content = []
            
            for cell_idx, cell in enumerate(notebook_data["cells"]):
                cell_type = cell.get("cell_type", "unknown")
                
                if cell_type == "markdown":
                    # Extract markdown content
                    source = cell.get("source", [])
                    if source:
                        markdown_text = "".join(source)
                        extracted_content.append(f"[Markdown Cell {cell_idx + 1}]\n{markdown_text}")
                
                elif cell_type == "code":
                    # Extract code content
                    source = cell.get("source", [])
                    if source:
                        code_text = "".join(source)
                        extracted_content.append(f"[Code Cell {cell_idx + 1}]\n```python\n{code_text}\n```")
                    
                    # Also extract outputs if present and meaningful
                    outputs = cell.get("outputs", [])
                    for output_idx, output in enumerate(outputs):
                        if output.get("output_type") == "execute_result":
                            data = output.get("data", {})
                            text_data = data.get("text/plain", [])
                            if text_data:
                                output_text = "".join(text_data)
                                extracted_content.append(f"[Code Cell {cell_idx + 1} Output {output_idx + 1}]\n{output_text}")
                
                elif cell_type == "raw":
                    # Extract raw content
                    source = cell.get("source", [])
                    if source:
                        raw_text = "".join(source)
                        extracted_content.append(f"[Raw Cell {cell_idx + 1}]\n{raw_text}")
            
            if not extracted_content:
                raise ValueError(f"No meaningful content found in notebook {ipynb_path.name}")
            
            combined_text = "\n\n".join(extracted_content)
            print(combined_text)
            print(f"Successfully extracted content from {len(extracted_content)} cells in {ipynb_path.name}")
            return combined_text
            
        except json.JSONDecodeError as e:
            raise Exception(f"Invalid JSON in notebook {ipynb_path.name}: {e}")
        except Exception as e:
            raise Exception(f"Failed to process notebook {ipynb_path.name}: {e}")
    
    def extract_pptx_content(self, pptx_path: Path) -> Dict[str, Any]:
        """
        Extract text and images from PowerPoint presentation (.pptx) file.
        
        Args:
            pptx_path: Path to the .pptx file
            
        Returns:
            Dictionary with 'text' (str) and 'images' (List[Tuple[int, Image.Image]])
            where images are (slide_number, PIL.Image) tuples
            
        Raises:
            Exception: If PowerPoint processing fails
        """
        if not PPTX_AVAILABLE:
            raise Exception("python-pptx not available. Install with: pip install python-pptx")
        
        try:
            prs = Presentation(pptx_path)
            slides_content = []
            images = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text_parts = []
                
                # Extract text and images from all shapes in the slide
                for shape in slide.shapes:
                    # Extract text
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:  # Only add non-empty text
                            slide_text_parts.append(text)
                    
                    # Extract images if PIL is available
                    if PIL_AVAILABLE:
                        try:
                            if hasattr(shape, "image") and shape.image is not None:
                                image_blob = shape.image.blob
                                pil_image = Image.open(BytesIO(image_blob)).convert("RGB")
                                images.append((slide_num, pil_image))
                                print(f"Extracted image from slide {slide_num}")
                        except Exception as e:
                            print(f"Warning: Could not extract image from slide {slide_num}: {e}")
                            continue
                
                # Combine all text from this slide
                if slide_text_parts:
                    slide_content = "\n".join(slide_text_parts)
                    slides_content.append(f"[Slide {slide_num}]\n{slide_content}")
            
            if not slides_content:
                raise ValueError(f"No text content found in PowerPoint {pptx_path.name}")
            
            combined_text = "\n\n".join(slides_content)
            print(f"Successfully extracted text from {len(slides_content)} slides and {len(images)} images from {pptx_path.name}")
            
            return {
                "text": combined_text,
                "images": images
            }
            
        except Exception as e:
            raise Exception(f"Failed to process PowerPoint {pptx_path.name}: {e}")

    def extract_tex_text(self, tex_path: Path) -> str:
        """Extract clean text from LaTeX, preserving sections and math."""
        try:
            with open(tex_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()

            # Remove comments
            content = re.sub(r'(?<!\\)%.*$', '', content, flags=re.MULTILINE)
            
            # Replace common commands with readable text
            replacements = {
                r'\\section\{([^}]+)\}': r'Section: \1',
                r'\\subsection\{([^}]+)\}': r'Subsection: \1',
                r'\\textbf\{([^}]+)\}': r'\1',
                r'\\textit\{([^}]+)\}': r'\1',
                r'\\emph\{([^}]+)\}': r'\1',
                r'\\cite\{[^}]+\}': '[Citation]',
                r'\\ref\{[^}]+\}': '[Reference]',
                r'\\label\{[^}]+\}': '',
                r'\\begin\{equation\}[\s\S]*?\\end\{equation\}': '[Equation]',
                r'\\begin\{align\}[\s\S]*?\\end\{align\}': '[Aligned Equations]',
                r'\$\$[\s\S]*?\$\$': '[Math Block]',
                r'\$[^$]+\$': '[Inline Math]',
            }
            
            for pattern, repl in replacements.items():
                content = re.sub(pattern, repl, content)
            
            # Clean up multiple newlines
            content = re.sub(r'\n\s*\n', '\n\n', content)
            return content.strip()
            
        except Exception as e:
            raise Exception(f"Failed to process .tex file {tex_path.name}: {e}")

    def extract_rmd_text(self, rmd_path: Path) -> str:
        """Extract structured text from R Markdown with code/prose separation."""
        try:
            with open(rmd_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()

            chunks = []
            lines = content.splitlines()
            i = 0
            in_code = False
            code_lang = ""
            current_chunk = []

            while i < len(lines):
                line = lines[i]

                # Detect YAML frontmatter
                if i == 0 and line.strip() == '---':
                    # Skip until next ---
                    while i < len(lines) and lines[i].strip() != '---':
                        i += 1
                    i += 1
                    continue

                # Detect code block start
                code_match = re.match(r'```\{?([a-zA-Z0-9,= ]*)\}?', line)
                if code_match:
                    if in_code:
                        # End previous code block
                        code_text = "\n".join(current_chunk).strip()
                        if code_text:
                            chunks.append(f"[R Code: {code_lang}]\n```r\n{code_text}\n```")
                        current_chunk = []
                        in_code = False
                    else:
                        # Start new code block
                        code_lang = code_match.group(1).strip().split(',')[0].replace('r', 'R').strip()
                        code_lang = code_lang or "R"
                        in_code = True
                        current_chunk = []
                    i += 1
                    continue

                if in_code:
                    current_chunk.append(line)
                else:
                    # Prose: clean inline R
                    line = re.sub(r'`r [^(]+`', '[Inline R]', line)  # `r mean(...)`
                    line = re.sub(r'```.*?```', '[Code]', line, flags=re.DOTALL)
                    if line.strip():
                        current_chunk.append(line)
                    else:
                        if current_chunk:
                            chunks.append(" ".join(current_chunk).strip())
                            current_chunk = []
                i += 1

            # Final prose chunk
            if not in_code and current_chunk:
                text = " ".join(current_chunk).strip()
                if text:
                    chunks.append(text)

            return "\n\n".join(chunks).strip()

        except Exception as e:
            raise Exception(f"Failed to process .rmd file {rmd_path.name}: {e}")
    
    def extract_html_text(self, html_path: Path) -> str:
        """
        Extract clean text from HTML files, stripping tags and preserving structure.
        
        Args:
            html_path: Path to the HTML file
            
        Returns:
            Extracted text content
        """
        if not BS4_AVAILABLE:
            # Fallback: basic regex extraction
            try:
                with open(html_path, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read()
                
                # Remove script and style tags and their content
                content = re.sub(r'<script[^>]*>.*?</script>', '', content, flags=re.DOTALL | re.IGNORECASE)
                content = re.sub(r'<style[^>]*>.*?</style>', '', content, flags=re.DOTALL | re.IGNORECASE)
                
                # Remove all HTML tags
                content = re.sub(r'<[^>]+>', '', content)
                
                # Clean up whitespace
                content = re.sub(r'\s+', ' ', content)
                content = re.sub(r'\n\s*\n', '\n\n', content)
                
                return content.strip()
                
            except Exception as e:
                raise Exception(f"Failed to process HTML file {html_path.name}: {e}")
        
        try:
            with open(html_path, "r", encoding="utf-8", errors="ignore") as f:
                html_content = f.read()
            
            # Parse HTML with BeautifulSoup
            soup = BeautifulSoup(html_content, 'lxml')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Get text and clean up whitespace
            text = soup.get_text()
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)
            
            return text.strip()
            
        except Exception as e:
            raise Exception(f"Failed to process HTML file {html_path.name}: {e}")
    
    def extract_text(self, file_path: Path, file_type: str) -> str:
        """
        Extract text from file based on its type.
        
        Args:
            file_path: Path to the file
            file_type: Type of file ('pdf', 'txt', 'ipynb', 'pptx', 'tex', 'rmd', 'html')
            
        Returns:
            Extracted text content
        """
        file_type_lower = file_type.lower()
        
        if "pdf" in file_type_lower:
            content = self.extract_pdf_content(file_path)
            return content["text"]
        elif "ipynb" in file_type_lower:
            return self.extract_ipynb_text(file_path)
        elif file_type_lower in ["pptx", "ppt"]:
            content = self.extract_pptx_content(file_path)
            return content["text"]
        elif file_type_lower == "tex":
            return self.extract_tex_text(file_path)
        elif file_type_lower == "rmd":
            return self.extract_rmd_text(file_path)
        elif file_type_lower == "html":
            return self.extract_html_text(file_path)
        else:
            return self.extract_text_file(file_path)

    
    
    def chunk_text(self, text: str) -> List[Document]:
        """
        Split text into chunks using RecursiveCharacterTextSplitter.
        
        Args:
            text: Input text to chunk
            
        Returns:
            List of Document objects with chunked text
        """
        if not text.strip():
            raise ValueError("Cannot chunk empty text")
        
        return self.text_splitter.create_documents([text])
    
    def create_metadata(self, doc_obj: Document, entry: Dict, source_file: Path) -> Dict:
        """
        Create metadata for a document chunk.
        
        Args:
            doc_obj: Langchain Document object
            entry: File entry from log
            source_file: Path to the source file
            
        Returns:
            Metadata dictionary for the chunk
        """
        # Extract page number from text if present (for PDFs)
        page_match = re.search(r"\[Page (\d+)\]", doc_obj.page_content)
        page_number = int(page_match.group(1)) if page_match else 1
        
        # Extract slide number from text if present (for PowerPoint files)
        slide_match = re.search(r"\[Slide (\d+)\]", doc_obj.page_content)
        slide_number = int(slide_match.group(1)) if slide_match else None
        
        # Extract cell information from Jupyter notebooks
        cell_type = None
        cell_number = None
        if "[Markdown Cell" in doc_obj.page_content:
            cell_match = re.search(r"\[Markdown Cell (\d+)\]", doc_obj.page_content)
            if cell_match:
                cell_type = "markdown"
                cell_number = int(cell_match.group(1))
        elif "[Code Cell" in doc_obj.page_content:
            cell_match = re.search(r"\[Code Cell (\d+)\]", doc_obj.page_content)
            if cell_match:
                cell_type = "code"
                cell_number = int(cell_match.group(1))
        elif "[Raw Cell" in doc_obj.page_content:
            cell_match = re.search(r"\[Raw Cell (\d+)\]", doc_obj.page_content)
            if cell_match:
                cell_type = "raw"
                cell_number = int(cell_match.group(1))
        
        metadata = {
            "vector_type": "chunk",  # Distinguishes from doc_summary vectors for layered RAG
            "id": entry.get("id"),
            "title": entry.get("title"),
            "course_id": entry.get("course_id"),
            "source": os.path.basename(source_file),
            "source_url": entry.get("source_url"),
            "type": entry.get("filetype"),
            "created_at": entry.get("created_at"),
            "modified_at": entry.get("modified_at"),
            "page_number": page_number,
            "size": entry.get("size"),
            "embedded_at": datetime.now().isoformat(),
            "text": doc_obj.page_content,  # Add the actual text content
        }
        
        # Add slide number for PowerPoint files
        if slide_number is not None:
            metadata["slide_number"] = slide_number
        
        # Add notebook-specific metadata
        if cell_type:
            metadata["cell_type"] = cell_type
            metadata["cell_number"] = cell_number
            metadata["notebook_chunk"] = True
        
        # Add module context if available (from module-first download)
        module_context = entry.get("module_context")
        if module_context:
            metadata["module_context"] = module_context
            metadata["module_id"] = module_context.get("module_id")
            metadata["module_position"] = module_context.get("position")
        
        return metadata


class EmbeddingClient:
    """
    Handles embedding model operations and vector generation.
    
    Manages SentenceTransformer model loading, text encoding,
    and batch processing for efficient embedding generation.
    """
    
    def __init__(self, model_name: str = "all-mpnet-base-v2", batch_size: int = 32):
        """
        Initialize the embedding client.
        
        Args:
            model_name: Name of the SentenceTransformer model
            batch_size: Batch size for embedding generation
        """
        self.model_name = model_name
        self.batch_size = batch_size
        self.model = None
        self.device = None
        self._initialize_model()
    
    def _initialize_model(self) -> None:
        """Initialize the SentenceTransformer model."""
        self.device = "cuda" if torch.cuda.is_available() else "cpu"
        print(f"Loading embedding model '{self.model_name}' on device: {self.device}")
        
        try:
            self.model = SentenceTransformer(self.model_name, device=self.device)
            print(f"Successfully loaded model on {self.device}")
        except Exception as e:
            raise Exception(f"Failed to load embedding model: {e}")
    
    def generate_embeddings(self, texts: List[str]) -> List[List[float]]:
        """
        Generate embeddings for a list of texts.
        
        Args:
            texts: List of text strings to embed
            
        Returns:
            List of embedding vectors
            
        Raises:
            Exception: If embedding generation fails
        """
        if not texts:
            raise ValueError("Cannot generate embeddings for empty text list")
        
        try:
            embeddings = self.model.encode(
                texts,
                batch_size=self.batch_size,
                convert_to_numpy=True,
                show_progress_bar=False,
                normalize_embeddings=True
            )
            return embeddings.tolist()
            
        except Exception as e:
            raise Exception(f"Failed to generate embeddings: {e}")
    
    def get_model_dimension(self) -> int:
        """
        Get the dimension of the embedding model.
        
        Returns:
            Embedding dimension
        """
        # all-mpnet-base-v2 outputs 768-dimensional embeddings
        return 768


class ImageEmbeddingClient:
    """
    Handles image embedding model operations using CLIP.
    
    Uses CLIP (Contrastive Language-Image Pre-training) model
    to generate embeddings for images that can be used for
    similarity search and retrieval.
    """
    
    def __init__(self, model_name: str = "clip-ViT-B-32", batch_size: int = 32):
        """
        Initialize the image embedding client.
        
        Args:
            model_name: Name of the CLIP model (e.g., "clip-ViT-B-32", "clip-ViT-L-14")
            batch_size: Batch size for embedding generation
        """
        self.model_name = model_name
        self.batch_size = batch_size
        self.model = None
        self.device = None
        self._initialize_model()
    
    def _initialize_model(self) -> None:
        """Initialize the CLIP model for image embeddings."""
        if not PIL_AVAILABLE:
            raise Exception("PIL/Pillow not available. Install with: pip install pillow")
        
        self.device = "cuda" if torch.cuda.is_available() else "cpu"
        print(f"Loading image embedding model '{self.model_name}' on device: {self.device}")
        
        try:
            self.model = SentenceTransformer(self.model_name, device=self.device)
            print(f"Successfully loaded image embedding model on {self.device}")
        except Exception as e:
            raise Exception(f"Failed to load image embedding model: {e}")
    
    def generate_image_embeddings(self, images: List['Image.Image']) -> List[List[float]]:
        """
        Generate embeddings for a list of PIL images.
        
        Args:
            images: List of PIL Image objects
            
        Returns:
            List of embedding vectors (512-dimensional for clip-ViT-B-32, 768 for clip-ViT-L-14)
            
        Raises:
            Exception: If embedding generation fails
        """
        if not images:
            raise ValueError("Cannot generate embeddings for empty image list")
        
        if not PIL_AVAILABLE:
            raise Exception("PIL/Pillow not available for image processing")
        
        try:
            embeddings = self.model.encode(
                images,
                batch_size=self.batch_size,
                convert_to_numpy=True,
                show_progress_bar=False,
                normalize_embeddings=True
            )
            return embeddings.tolist()
            
        except Exception as e:
            raise Exception(f"Failed to generate image embeddings: {e}")
    
    def get_model_dimension(self) -> int:
        """
        Get the dimension of the image embedding model.
        
        Returns:
            Embedding dimension (512 for clip-ViT-B-32, 768 for clip-ViT-L-14)
        """
        # Default to clip-ViT-B-32 which outputs 512-dimensional embeddings
        # If using clip-ViT-L-14, it outputs 768-dimensional embeddings
        if "ViT-L" in self.model_name:
            return 768
        else:
            return 512


class VectorStore:
    """
    Manages Pinecone vector database operations.
    
    Handles index creation, vector upsert operations,
    and connection management with the Pinecone service.
    """
    
    def __init__(self, api_key: str, index_name: str = "course-documents", dimension: int = 768):
        """
        Initialize the vector store.
        
        Args:
            api_key: Pinecone API key
            index_name: Name of the Pinecone index
            dimension: Dimension of the vectors (default 768 for text embeddings)
        """
        self.api_key = api_key
        self.index_name = index_name
        self.dimension = dimension
        self.pinecone_client = None
        self.index = None
        self._initialize_client()
    
    def _initialize_client(self) -> None:
        """Initialize Pinecone client and ensure index exists."""
        try:
            print("Connecting to Pinecone...")
            self.pinecone_client = Pinecone(api_key=self.api_key)
            self.index = self._ensure_index_exists()
            print(f"Connected to Pinecone index: {self.index_name}")
            
        except Exception as e:
            raise Exception(f"Failed to initialize Pinecone client: {e}")
    
    def _ensure_index_exists(self) -> 'Pinecone.Index':
        """
        Create the index if it doesn't exist.
        
        Returns:
            Pinecone index object
        """
        existing_indexes = self.pinecone_client.list_indexes().names()
        
        if self.index_name not in existing_indexes:
            print(f"Creating Pinecone index: {self.index_name} with dimension {self.dimension}")
            self.pinecone_client.create_index(
                name=self.index_name,
                dimension=self.dimension,
                metric="cosine",
                spec=ServerlessSpec(cloud="aws", region="us-east-1")
            )
            print(f"Created index: {self.index_name}")
        
        return self.pinecone_client.Index(self.index_name)
    
    def upsert_vectors(self, vectors: List[Dict]) -> None:
        """
        Upsert vectors to the Pinecone index.
        
        Args:
            vectors: List of vector dictionaries with 'id', 'values', and 'metadata'
            
        Raises:
            Exception: If upsert operation fails
        """
        if not vectors:
            return
        
        try:
            self.index.upsert(vectors=vectors)
        except Exception as e:
            raise Exception(f"Failed to upsert vectors to Pinecone: {e}")
    
    def check_vectors_exist(self, document_id: str) -> bool:
        """
        Check if vectors for a document already exist in the index.
        
        Uses fetch with sample vector IDs that would be generated for this document.
        
        Args:
            document_id: The document ID to check for
            
        Returns:
            True if vectors exist for this document, False otherwise
        """
        try:
            stats = self.index.describe_index_stats()
            
            # If index is empty, no vectors exist
            if stats.total_vector_count == 0:
                return False
            
            # Check if any vectors with the document ID prefix exist
            # We'll check a few sample IDs that would be generated for this document
            # Format: {document_id}_{page_number}_{chunk_index}
            sample_ids = [
                f"{document_id}_1_0",  # First chunk, first page
                f"{document_id}_1_1",  # Second chunk, first page
                f"{document_id}_2_0",  # First chunk, second page
            ]
            
            try:
                fetch_result = self.index.fetch(ids=sample_ids)
                if fetch_result and fetch_result.vectors:
                    # If any of the sample IDs exist, the document is likely already embedded
                    return len(fetch_result.vectors) > 0
            except Exception as e:
                # If fetch fails, assume vectors don't exist to be safe
                return False
            
            return False
                
        except Exception as e:
            print(f"Warning: Could not check if vectors exist for document {document_id}: {e}")
            # On error, return False to allow processing (safer than blocking)
            return False
    
    def get_index_stats(self) -> Dict:
        """
        Get statistics about the Pinecone index.
        
        Returns:
            Dictionary with index statistics
        """
        try:
            return self.index.describe_index_stats()
        except Exception as e:
            print(f"Warning: Could not get index stats: {e}")
            return {}


class TextStorage:
    """
    Manages storage of raw text chunks for inspection and debugging.
    
    Creates organized text files that correspond to the embedded chunks,
    allowing for easy inspection of what content was processed.
    """
    
    def __init__(self, storage_root: Path):
        """
        Initialize the text storage manager.
        
        Args:
            storage_root: Root directory for storing raw text files
        """
        self.storage_root = Path(storage_root)
        self.text_chunks_dir = self.storage_root / "text_chunks"
        self._ensure_storage_dir()
    
    def _ensure_storage_dir(self) -> None:
        """Ensure the text chunks storage directory exists."""
        self.text_chunks_dir.mkdir(parents=True, exist_ok=True)
    
    def save_text_chunks(self, entry: Dict, texts: List[str]) -> None:
        """
        Save raw text chunks to organized files for inspection.
        
        Args:
            entry: Log entry for the document
            texts: List of text chunks
        """
        try:
            # Create a directory for this document
            doc_id = entry.get('id', 'unknown')
            course_id = entry.get('course_id', 'unknown')
            title = entry.get('title', 'unknown').replace('/', '_').replace('\\', '_')
            
            doc_dir = self.text_chunks_dir / f"course_{course_id}" / f"doc_{doc_id}_{title}"
            doc_dir.mkdir(parents=True, exist_ok=True)
            
            # Save each chunk as a separate file
            for i, text in enumerate(texts):
                # Create filename with chunk info
                chunk_filename = f"chunk_{i:03d}.txt"
                chunk_path = doc_dir / chunk_filename
                
                # Write the text chunk with simple header
                with open(chunk_path, 'w', encoding='utf-8') as f:
                    f.write(f"# Chunk {i} from {title}\n")
                    f.write(f"# Course ID: {course_id}\n")
                    f.write(f"# Document ID: {doc_id}\n")
                    f.write(f"# Timestamp: {datetime.now().isoformat()}\n")
                    f.write("# " + "="*60 + "\n\n")
                    f.write(text)
            
            # Save a summary file for the document
            summary_path = doc_dir / "summary.txt"
            with open(summary_path, 'w', encoding='utf-8') as f:
                f.write(f"Document Processing Summary\n")
                f.write("="*40 + "\n\n")
                f.write(f"Title: {title}\n")
                f.write(f"Document ID: {doc_id}\n")
                f.write(f"Course ID: {course_id}\n")
                f.write(f"Total Chunks: {len(texts)}\n")
                f.write(f"Processing Time: {datetime.now().isoformat()}\n\n")
                
                f.write("Chunk Overview:\n")
                f.write("-" * 20 + "\n")
                for i in range(len(texts)):
                    f.write(f"Chunk {i}: chunk_{i:03d}.txt\n")
            
            print(f"Saved {len(texts)} text chunks to {doc_dir}")
            
        except Exception as e:
            print(f"Warning: Failed to save text chunks for {entry.get('title', 'unknown')}: {e}")
    
    def save_images(self, entry: Dict, images: List[Tuple[int, 'Image.Image']], location_type: str = "page") -> List[str]:
        """
        Save extracted images to organized files for inspection.
        
        Args:
            entry: Log entry for the document
            images: List of (location_number, PIL.Image) tuples
            location_type: Type of location ("page" for PDFs, "slide" for PPTX)
            
        Returns:
            List of local file paths where images were saved (in same order as input images)
        """
        if not PIL_AVAILABLE:
            print("Warning: PIL not available, cannot save images")
            return []
            
        saved_paths = []
        try:
            # Create a directory for this document
            doc_id = entry.get('id', 'unknown')
            course_id = entry.get('course_id', 'unknown')
            title = entry.get('title', 'unknown').replace('/', '_').replace('\\', '_')
            
            doc_dir = self.text_chunks_dir / f"course_{course_id}" / f"doc_{doc_id}_{title}"
            images_dir = doc_dir / "images"
            images_dir.mkdir(parents=True, exist_ok=True)
            
            # Save each image with proper indexing
            location_counts = {}  # Track how many images per location
            for location_num, pil_image in images:
                # Count images for this location
                if location_num not in location_counts:
                    location_counts[location_num] = 0
                location_counts[location_num] += 1
                
                # Create filename with location info
                image_filename = f"{location_type}_{location_num:03d}_image_{location_counts[location_num]:02d}.png"
                image_path = images_dir / image_filename
                
                # Save the image
                pil_image.save(image_path, "PNG")
                
                # Store the path as a string (relative to storage_root/ref_root for portability)
                # Path will be like: text_chunks/course_{course_id}/doc_{doc_id}_{title}/images/page_001_image_01.png
                relative_path = image_path.relative_to(self.storage_root)
                saved_paths.append(str(relative_path))
            
            # Update summary file with image info
            summary_path = doc_dir / "summary.txt"
            if summary_path.exists():
                with open(summary_path, 'a', encoding='utf-8') as f:
                    f.write(f"\nTotal Images: {len(images)}\n")
                    f.write(f"Images saved to: images/\n")
            
            print(f"Saved {len(images)} images to {images_dir}")
            
        except Exception as e:
            print(f"Warning: Failed to save images for {entry.get('title', 'unknown')}: {e}")
        
        return saved_paths


class LogEntryManager:
    """
    Manages reading and filtering of log entries from the files log.
    
    Handles JSON parsing, file filtering, and validation
    of log entries for processing.
    """
    
    @staticmethod
    def load_log_entries(log_path: Path) -> List[Dict]:
        """
        Load log entries from a JSONL file.
        
        Args:
            log_path: Path to the log file
            
        Returns:
            List of log entry dictionaries
            
        Raises:
            Exception: If file reading fails
        """
        if not log_path.exists():
            raise FileNotFoundError(f"Log file not found: {log_path}")
        
        entries = []
        
        try:
            with open(log_path, "r", encoding="utf-8") as f:
                for line_num, line in enumerate(f, 1):
                    line = line.strip()
                    if not line:
                        continue

                    try:
                        entry = json.loads(line)
                        entries.append(entry)
                    except json.JSONDecodeError as e:
                        print(f"Warning: Invalid JSON at line {line_num}: {e}")
                        continue
                        
        except Exception as e:
            raise Exception(f"Failed to read log file {log_path}: {e}")
        
        return entries
    
    @staticmethod
    def filter_eligible_files(entries: List[Dict]) -> List[Dict]:
        """
        Filter log entries to only include eligible files (PDF, TXT, IPYNB, PPTX, TEX, and RMD).
        
        Args:
            entries: List of all log entries
            
        Returns:
            List of filtered entries for eligible files
        """
        eligible_files = []
        supported_types = {"pdf", "txt", "ipynb", "pptx", "ppt", "tex", "rmd", "html"}
        
        for entry in entries:
            try:
                # Check if file exists
                local_path = Path(entry.get("local_path", ""))
                if not local_path.exists():
                    print(f"Warning: File not found: {local_path}")
                    continue
                
                # Check file type
                title = entry.get("title", "")
                file_type = str(local_path).split(".")[-1].lower() if "." in str(local_path) else ""

                
                if file_type in supported_types:
                    eligible_files.append(entry)
                else:
                    print(f"Skipping unsupported file type: {title} - {file_type}")
                    
            except Exception as e:
                print(f"Warning: Error processing entry {entry.get('id', 'unknown')}: {e}")
                continue

        return eligible_files


class EmbeddingPipeline:
    """
    Main orchestrator for the embedding pipeline.
    
    Coordinates document processing, embedding generation,
    and vector storage operations with comprehensive reporting.
    """
    
    def __init__(self, env_path: Path, ref_root: Path, 
                 chunk_size: int = 500, chunk_overlap: int = 100,
                 batch_size: int = 64, save_raw_text: bool = True):
        """
        Initialize the embedding pipeline.
        
        Args:
            env_path: Path to .env file with API credentials
            ref_root: Root directory for reference files and logs
            chunk_size: Size of text chunks
            chunk_overlap: Overlap between chunks
            batch_size: Batch size for operations
            save_raw_text: Whether to save raw text chunks for inspection
        """
        self.env_path = env_path
        self.ref_root = ref_root
        self.log_file = ref_root / "files.jsonl"
        self.report_file = ref_root / "embedding_report.json"
        self.save_raw_text = save_raw_text
        
        # Initialize components
        self.document_processor = DocumentProcessor(chunk_size, chunk_overlap)
        self.embedding_client = EmbeddingClient(batch_size=batch_size)
        self.vector_store = VectorStore(self._load_pinecone_key(), "course-documents")
        
        # Initialize image embedding components
        try:
            self.image_embedding_client = ImageEmbeddingClient(batch_size=batch_size)
            image_dimension = self.image_embedding_client.get_model_dimension()
            self.image_vector_store = VectorStore(
                self._load_pinecone_key(), 
                "image-embeddings",
                dimension=image_dimension
            )
            print(f"Image embedding system initialized with dimension {image_dimension}")
        except Exception as e:
            print(f"Warning: Could not initialize image embedding system: {e}")
            print("Image embedding will be skipped. Install required packages: pip install sentence-transformers pillow")
            self.image_embedding_client = None
            self.image_vector_store = None
        
        # Initialize text storage if enabled
        if self.save_raw_text:
            self.text_storage = TextStorage(ref_root)
    
    def _load_pinecone_key(self) -> str:
        """
        Load Pinecone API key from environment.
        
        Returns:
            Pinecone API key
            
        Raises:
            ValueError: If API key is missing
        """
        load_dotenv(self.env_path)
        api_key = os.getenv("PINECONE_API_KEY")

        
        if not api_key:
            raise ValueError("Missing PINECONE_API_KEY in .env file")
        
        return api_key
    
    def process_document(self, entry: Dict) -> Tuple[List[Dict], int]:
        """
        Process a single document through the embedding pipeline.
        
        Args:
            entry: Log entry for the document
            
        Returns:
            Tuple of (vector_data, chunk_count)
        """
        file_path = Path(entry["local_path"])
        # file_type = entry.get("filetype", "").lower()
        file_type = entry["local_path"].split(".")[-1].lower() if "." in entry["local_path"] else ""
        
        # Extract text and images from document (for PDF/PPTX, extract both)
        text = None
        images = []
        location_type = "page"
        
        if file_type in ["pdf"]:
            # Extract both text and images from PDF
            content = self.document_processor.extract_pdf_content(file_path)
            text = content["text"]
            images = content["images"]
            location_type = "page"
        elif file_type in ["pptx", "ppt"]:
            # Extract both text and images from PPTX
            content = self.document_processor.extract_pptx_content(file_path)
            text = content["text"]
            images = content["images"]
            location_type = "slide"
        else:
            # For other file types, just extract text
            text = self.document_processor.extract_text(file_path, file_type)
        
        # Save extracted images if any were found
        image_paths = []
        if self.save_raw_text and self.text_storage and images:
            image_paths = self.text_storage.save_images(entry, images, location_type)
        
        # Chunk the text
        docs = self.document_processor.chunk_text(text)
        texts = [doc.page_content for doc in docs]
        
        print(f"Chunked {len(texts)} chunks from {file_path.name}")
        if images:
            print(f"Extracted {len(images)} images from {file_path.name}")
        
        # Save raw text chunks for inspection if enabled
        if self.save_raw_text and self.text_storage:
            self.text_storage.save_text_chunks(entry, texts)

        # Generate embeddings
        embeddings = self.embedding_client.generate_embeddings(texts)
        
        # Prepare vectors for upsert
        vectors = []
        for i, (doc_obj, emb) in enumerate(zip(docs, embeddings)):
            metadata = self.document_processor.create_metadata(doc_obj, entry, file_path)
            
            # Clean metadata to remove None/null values (Pinecone doesn't accept them)
            cleaned_metadata = self._clean_metadata_for_pinecone(metadata)
            
            # Create vector ID based on available location info
            # For PowerPoint files, use slide_number if available, otherwise fall back to page_number
            location_id = cleaned_metadata.get('slide_number', cleaned_metadata.get('page_number', 1))
            vector_id = f"{entry['id']}_{location_id}_{i}"
            print(vector_id)
            
            vectors.append({
                "id": vector_id,
                "values": emb,
                "metadata": cleaned_metadata,
            })

        # Build location-to-text mapping for context extraction
        location_to_text_chunks = {}
        for doc_obj in docs:
            # Extract page/slide number from chunk metadata
            page_match = re.search(r"\[Page (\d+)\]", doc_obj.page_content)
            slide_match = re.search(r"\[Slide (\d+)\]", doc_obj.page_content)
            
            if location_type == "page" and page_match:
                page_num = int(page_match.group(1))
                if page_num not in location_to_text_chunks:
                    location_to_text_chunks[page_num] = []
                location_to_text_chunks[page_num].append(doc_obj.page_content)
            elif location_type == "slide" and slide_match:
                slide_num = int(slide_match.group(1))
                if slide_num not in location_to_text_chunks:
                    location_to_text_chunks[slide_num] = []
                location_to_text_chunks[slide_num].append(doc_obj.page_content)
        
        # Embed and store images if available
        image_vectors = []
        if images and self.image_embedding_client and self.image_vector_store:
            try:
                print(f"Generating embeddings for {len(images)} images...")
                # Extract just the PIL images from the tuples
                pil_images = [img for _, img in images]
                image_embeddings = self.image_embedding_client.generate_image_embeddings(pil_images)
                
                # Create vectors for images
                for idx, ((location_num, pil_image), emb) in enumerate(zip(images, image_embeddings)):
                    # Create metadata for image
                    image_metadata = {
                        "id": entry.get("id"),
                        "title": entry.get("title"),
                        "course_id": entry.get("course_id"),
                        "source": os.path.basename(file_path),
                        "source_url": entry.get("source_url"),
                        "type": entry.get("filetype"),
                        "modality": "image",
                        "embedded_at": datetime.now().isoformat(),
                        "size": entry.get("size"),
                    }
                    
                    # Add location-specific metadata (page_number or slide_number)
                    if location_type == "page":
                        image_metadata["page_number"] = location_num
                    elif location_type == "slide":
                        image_metadata["slide_number"] = location_num
                    
                    # Add local path to saved image (relative to ref_root)
                    # Path format: text_chunks/course_{course_id}/doc_{doc_id}_{title}/images/page_001_image_01.png
                    # Full path can be reconstructed as: {ref_root}/{image_local_path}
                    if idx < len(image_paths) and image_paths[idx]:
                        image_metadata["image_local_path"] = image_paths[idx]
                    
                    # Extract contextual text from the same location (page/slide)
                    # This enables text-based queries to find relevant images in mid/late fusion RAG
                    context_text = ""
                    if location_num in location_to_text_chunks:
                        # Get all text chunks from the same page/slide
                        location_chunks = location_to_text_chunks[location_num]
                        # Remove location markers for cleaner context
                        cleaned_chunks = []
                        for chunk in location_chunks:
                            # Remove [Page X] or [Slide X] markers
                            cleaned = re.sub(rf"\[{location_type.title()} {location_num}\]\n?", "", chunk)
                            cleaned_chunks.append(cleaned.strip())
                        # Combine chunks with newlines
                        context_text = "\n\n".join(cleaned_chunks)
                    
                    # Add contextual text to metadata (for text-based retrieval and LLM context)
                    if context_text:
                        # Truncate if too long (Pinecone metadata has limits, keep it reasonable)
                        max_context_length = 2000  # characters
                        if len(context_text) > max_context_length:
                            context_text = context_text[:max_context_length] + "..."
                        image_metadata["context_text"] = context_text
                        image_metadata["has_context"] = True
                    else:
                        image_metadata["has_context"] = False
                    
                    # Also add document title and section info as additional context
                    doc_title = entry.get("title", "")
                    if doc_title:
                        image_metadata["document_title"] = doc_title
                    
                    # Clean metadata
                    cleaned_image_metadata = self._clean_metadata_for_pinecone(image_metadata)
                    
                    # Create image vector ID
                    image_vector_id = f"{entry['id']}_{location_type}_{location_num}_img{idx}"
                    
                    image_vectors.append({
                        "id": image_vector_id,
                        "values": emb,
                        "metadata": cleaned_image_metadata,
                    })
                
                print(f"Created {len(image_vectors)} image vectors")
                
                # Upsert image vectors to image-embeddings index
                if image_vectors:
                    print(f"Upserting {len(image_vectors)} image vectors to image-embeddings index...")
                    self._batch_upsert_image_vectors(image_vectors)
                    print("Image vectors upserted successfully")
                    
            except Exception as e:
                print(f"Warning: Failed to embed and store images: {e}")
                import traceback
                traceback.print_exc()

        # Layered RAG: one document-level vector = mean of chunk embeddings (no LLM, same index)
        if vectors:
            try:
                dim = len(vectors[0]["values"])
                mean_embedding = [
                    sum(v["values"][i] for v in vectors) / len(vectors)
                    for i in range(dim)
                ]
                doc_preview = (docs[0].page_content[:500] + "...") if docs else entry.get("title", "")
                summary_meta = {
                    "vector_type": "doc_summary",
                    "id": entry.get("id"),
                    "title": entry.get("title"),
                    "course_id": entry.get("course_id"),
                    "text": doc_preview,
                    "source": entry.get("title", ""),
                    "embedded_at": datetime.now().isoformat(),
                }
                summary_meta_clean = self._clean_metadata_for_pinecone(summary_meta)
                vectors.append({
                    "id": f"summary_{entry['id']}",
                    "values": mean_embedding,
                    "metadata": summary_meta_clean,
                })
            except Exception as e:
                print(f"Warning: Could not add doc summary vector: {e}")

        return vectors, len(texts)
    
    def _clean_metadata_for_pinecone(self, metadata: Dict) -> Dict:
        """
        Clean metadata to remove None/null values and flatten nested dicts for Pinecone.
        
        Pinecone metadata requirements:
        - Must be string, number, boolean, or list of strings
        - Cannot contain None/null values
        - Cannot contain nested dictionaries (must be flattened)
        
        Args:
            metadata: Raw metadata dictionary
            
        Returns:
            Cleaned metadata dictionary with flattened structure
        """
        cleaned = {}
        
        for key, value in metadata.items():
            # Skip None/null values
            if value is None:
                continue
            
            # Handle nested dictionaries (like module_context) - flatten them
            if isinstance(value, dict):
                # Flatten nested dict with prefix
                for nested_key, nested_value in value.items():
                    if nested_value is not None:
                        # Use key_nestedkey as the flattened key
                        flat_key = f"{key}_{nested_key}"
                        cleaned[flat_key] = nested_value
                continue
            
            # Handle lists
            if isinstance(value, list):
                # Filter out None values from lists
                cleaned_list = [item for item in value if item is not None]
                if cleaned_list:
                    cleaned[key] = cleaned_list
                continue
            
            # Handle strings, numbers, booleans
            if isinstance(value, (str, int, float, bool)):
                cleaned[key] = value
            # Convert other types to string
            elif value is not None:
                cleaned[key] = str(value)
        
        return cleaned
    
    def run_pipeline(self) -> None:
        """
        Run the complete embedding pipeline.
        
        Processes all eligible files, generates embeddings,
        and stores them in Pinecone with comprehensive reporting.
        """
        print("Starting embedding pipeline...")
        
        try:
            # Load and filter log entries
            print("Reading log entries...")
            all_entries = LogEntryManager.load_log_entries(self.log_file)
            eligible_files = LogEntryManager.filter_eligible_files(all_entries)
            
            print(f"Found {len(eligible_files)} eligible files to embed")
            
            if not eligible_files:
                print("No eligible files found. Pipeline complete.")
                return
            
            total_chunks = 0
            reports = []
            skipped_count = 0
            
            # Process each file
            for entry in eligible_files:
                try:
                    document_id = entry.get("id")
                    file_name = Path(entry["local_path"]).name
                    file_type = entry["local_path"].split(".")[-1].lower() if "." in entry["local_path"] else ""
                    
                    # Check if vectors already exist for this document
                    # Skip the check for PDF files (allow reprocessing)
                    if file_type not in ["pdf", "pptx", "ppt"] and document_id and self.vector_store.check_vectors_exist(document_id):
                        print(f"Skipping {file_name} - vectors already exist in index")
                        skipped_count += 1
                        reports.append({
                            "file": file_name,
                            "chunks": 0,
                            "course": entry["course_id"],
                            "timestamp": datetime.now().isoformat(),
                            "status": "skipped (already embedded)"
                        })
                        continue
                    
                    vectors, chunk_count = self.process_document(entry)
                    
                    # Batch upsert vectors
                    print("Upserting vectors...")
                    self._batch_upsert_vectors(vectors)
                    print("Vectors upserted successfully")
                    total_chunks += chunk_count
                    
                    # Create report entry
                    reports.append({
                        "file": file_name,
                        "chunks": chunk_count,
                        "course": entry["course_id"],
                        "timestamp": datetime.now().isoformat(),
                        "status": "processed"
                    })

                    print(f"Successfully processed {file_name}")
                    
                except Exception as e:
                    print(f"Failed to process {entry.get('title', 'unknown')}: {e}")
                    continue
            
            # Save processing report
            self._save_report(reports)
            
            # Regenerate database overview (TOC) for RAG agent
            try:
                from database_overview import generate_database_overview
                generate_database_overview(self.ref_root, use_embedding_report=True)
                print(f"Database overview (TOC) updated: {self.ref_root / 'database_overview.txt'}")
            except Exception as e:
                print(f"Warning: Could not update database overview: {e}")
            
            print(f"\nEmbedding pipeline completed successfully!")
            processed_count = len(eligible_files) - skipped_count
            print(f"Processed {processed_count} files, skipped {skipped_count} files (already embedded)")
            print(f"Embedded {total_chunks} chunks")
            print(f"Report saved to: {self.report_file}")
            
        except Exception as e:
            print(f"Embedding pipeline failed: {e}")
            raise
    
    def _batch_upsert_vectors(self, vectors: List[Dict]) -> None:
        """
        Upsert vectors in batches for efficiency.
        
        Args:
            vectors: List of vectors to upsert
        """
        batch_size = self.embedding_client.batch_size
        
        for i in range(0, len(vectors), batch_size):
            batch = vectors[i:i + batch_size]
            self.vector_store.upsert_vectors(batch)
    
    def _batch_upsert_image_vectors(self, image_vectors: List[Dict]) -> None:
        """
        Upsert image vectors in batches for efficiency.
        
        Args:
            image_vectors: List of image vectors to upsert
        """
        if not self.image_vector_store:
            return
        
        batch_size = self.image_embedding_client.batch_size if self.image_embedding_client else 32
        
        for i in range(0, len(image_vectors), batch_size):
            batch = image_vectors[i:i + batch_size]
            self.image_vector_store.upsert_vectors(batch)
    
    def _save_report(self, reports: List[Dict]) -> None:
        """
        Save processing report to JSON file.
        
        Args:
            reports: List of processing reports
        """
        try:
            with open(self.report_file, "w", encoding="utf-8") as f:
                json.dump(reports, f, indent=2)
        except Exception as e:
            print(f"Warning: Could not save report to {self.report_file}: {e}")


def main():
    """
    Main entry point for the embedding pipeline.
    
    Sets up configuration and runs the complete embedding process.
    """
    # Configuration paths
    PROJECT_ROOT = Path(__file__).resolve().parents[0]
    REF_ROOT = PROJECT_ROOT / "ref_data"
    ENV_PATH = REF_ROOT / ".env"
    
    print("Canvas Document Embedding Pipeline")
    print(f"Loading configuration from {ENV_PATH}")
    print(f"Using reference data from {REF_ROOT}")
    print(f"Raw text chunks will be saved to: {REF_ROOT / 'text_chunks'}")
    
    try:
        # Initialize and run pipeline
        pipeline = EmbeddingPipeline(ENV_PATH, REF_ROOT)
        pipeline.run_pipeline()
        
        return 0
        
    except Exception as e:
        print(f"Pipeline execution failed: {e}")
        return 1


if __name__ == "__main__":
    exit(main())
